#!/usr/bin/env python3
"""
docs/_build_pptx.py
-------------------
Build ``docs/presentation.pptx`` — the 15-minute defense deck for
the UIR S8 Integrated Project (Product Review Intelligence).

Design system
    Background   #0D1B2A (dark navy)         — committed throughout
    Primary text #FFFFFF (white)             — bodies and titles
    Accent       #C9A84C (muted gold)        — motif bar, stat callouts, numbers
    Muted        #8A9BB4 (cool slate)        — captions, footers
    Card fill    #152a3f (navy +1 step)      — content blocks

Visual motif (repeated every slide)
    - Thin gold bar on the left edge (0.12" wide, full height)
    - "Product Review Intelligence"  footer bottom-left in slate
    - Slide number "N / 11" bottom-right in gold

Typography
    Header: Georgia           — gives the talk a report-like authority
    Body:   Calibri           — clean on projectors

Usage
    /opt/homebrew/bin/python3.10 docs/_build_pptx.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
ROOT = Path(__file__).resolve().parent
OUT = ROOT / "presentation.pptx"

# Colour tokens
NAVY = RGBColor(0x0D, 0x1B, 0x2A)
CARD = RGBColor(0x15, 0x2A, 0x3F)
CARD_DARKER = RGBColor(0x0F, 0x22, 0x35)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xC9, 0xA8, 0x4C)
GOLD_DIM = RGBColor(0x8B, 0x75, 0x33)
SLATE = RGBColor(0x8A, 0x9B, 0xB4)
RED = RGBColor(0xD9, 0x6E, 0x6E)
GREEN = RGBColor(0x7F, 0xB0, 0x7F)

# Fonts
H_FONT = "Georgia"
B_FONT = "Calibri"

TOTAL_SLIDES = 11
AUTHORS = (
    "Mohamed Anouar FARESS   •   Mohammed Yassine BEN MEZIANE\n"
    "Imad Eddine EL KHORASSANI   •   Anys LOUZAL"
)

# 16:9 slide: 13.333 x 7.5 inches
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------
def solid_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # no border by default


def add_rect(slide, x, y, w, h, color: RGBColor):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    solid_fill(shp, color)
    return shp


def add_round_rect(slide, x, y, w, h, color: RGBColor, radius: float = 0.04):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = radius
    solid_fill(shp, color)
    return shp


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text: str,
    *,
    font: str = B_FONT,
    size: int = 14,
    color: RGBColor = WHITE,
    bold: bool = False,
    italic: bool = False,
    align: str = "left",
    anchor: str = "top",
    line_spacing: float | None = None,
):
    """Add a plain text box.

    ``text`` may contain newlines — each becomes a paragraph. For richer
    multi-run formatting use :func:`add_rich_text`.
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }[anchor]

    align_enum = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }[align]

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_enum
        if line_spacing is not None:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_bullets(
    slide,
    x,
    y,
    w,
    h,
    items: list[str],
    *,
    size: int = 16,
    color: RGBColor = WHITE,
    gap: float = 1.15,
    bullet_color: RGBColor = GOLD,
):
    """Add a list with gold-square bullets.

    Uses a Unicode ▪ glyph coloured gold, followed by a space and then the
    text. This avoids PowerPoint's default bullet XML (which is fiddly in
    python-pptx) while still looking intentional.
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = gap
        p.space_after = Pt(4)

        bullet_run = p.add_run()
        bullet_run.text = "▪  "
        bullet_run.font.name = B_FONT
        bullet_run.font.size = Pt(size)
        bullet_run.font.color.rgb = bullet_color
        bullet_run.font.bold = True

        body_run = p.add_run()
        body_run.text = item
        body_run.font.name = B_FONT
        body_run.font.size = Pt(size)
        body_run.font.color.rgb = color
    return tb


# ---------------------------------------------------------------------------
# Chrome (per-slide frame)
# ---------------------------------------------------------------------------
def apply_chrome(slide, slide_number: int, *, show_footer: bool = True) -> None:
    """Paint navy background, gold motif bar, and footer on every slide."""
    # Background fill
    bg = add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    # Send to back (new shapes are inserted at the front)
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)

    # Gold motif bar — left edge, full height
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, GOLD)

    if show_footer:
        # Footer label (bottom-left)
        add_text(
            slide,
            Inches(0.45),
            Inches(7.1),
            Inches(6),
            Inches(0.3),
            "PRODUCT REVIEW INTELLIGENCE  ·  UIR ESIN IA S8",
            size=9,
            color=SLATE,
        )
        # Slide number (bottom-right)
        add_text(
            slide,
            Inches(11.8),
            Inches(7.1),
            Inches(1.3),
            Inches(0.3),
            f"{slide_number:02d} / {TOTAL_SLIDES:02d}",
            size=10,
            color=GOLD,
            align="right",
            bold=True,
        )


def add_slide(prs: Presentation) -> "slide":
    blank = prs.slide_layouts[6]  # truly blank
    return prs.slides.add_slide(blank)


def title_block(slide, title: str, eyebrow: str | None = None) -> None:
    """Standard title block used on content slides."""
    if eyebrow:
        add_text(
            slide,
            Inches(0.6),
            Inches(0.55),
            Inches(10),
            Inches(0.3),
            eyebrow.upper(),
            size=11,
            color=GOLD,
            bold=True,
        )
    add_text(
        slide,
        Inches(0.6),
        Inches(0.85),
        Inches(12),
        Inches(0.9),
        title,
        font=H_FONT,
        size=32,
        color=WHITE,
        bold=True,
    )


# ===========================================================================
# Slide builders
# ===========================================================================
def slide_1_title(prs: Presentation) -> None:
    s = add_slide(prs)
    apply_chrome(s, 1, show_footer=False)

    # Oversized muted "01" in the top-right — design motif
    add_text(
        s,
        Inches(11.4),
        Inches(0.4),
        Inches(1.6),
        Inches(1.2),
        "01",
        font=H_FONT,
        size=60,
        color=CARD,
        bold=True,
        align="right",
    )

    # Small eyebrow above the title
    add_text(
        s,
        Inches(1.0),
        Inches(1.6),
        Inches(11),
        Inches(0.4),
        "UIR  ·  ESIN  ·  IA S8  ·  INTEGRATED PROJECT",
        size=13,
        color=GOLD,
        bold=True,
    )

    # Gold underline accent (short, left-aligned — NOT a full-width accent line)
    add_rect(s, Inches(1.0), Inches(2.05), Inches(0.6), Inches(0.04), GOLD)

    # Main title
    add_text(
        s,
        Inches(1.0),
        Inches(2.3),
        Inches(11.3),
        Inches(1.3),
        "Product Review Intelligence",
        font=H_FONT,
        size=54,
        color=WHITE,
        bold=True,
    )

    # Subtitle
    add_text(
        s,
        Inches(1.0),
        Inches(3.7),
        Inches(11.3),
        Inches(1.0),
        "A Multi-Agent AI System for Sentiment Analysis\n& Market Research",
        font=H_FONT,
        size=22,
        color=SLATE,
        italic=True,
        line_spacing=1.2,
    )

    # Authors card — bottom, on CARD fill
    add_rect(s, Inches(1.0), Inches(5.5), Inches(11.3), Inches(0.08), GOLD)
    add_text(
        s,
        Inches(1.0),
        Inches(5.65),
        Inches(11.3),
        Inches(0.35),
        "AUTHORS",
        size=10,
        color=GOLD,
        bold=True,
    )
    add_text(
        s,
        Inches(1.0),
        Inches(6.0),
        Inches(11.3),
        Inches(0.8),
        AUTHORS,
        size=15,
        color=WHITE,
        line_spacing=1.3,
    )
    add_text(
        s,
        Inches(1.0),
        Inches(6.9),
        Inches(11.3),
        Inches(0.3),
        "APRIL 2026",
        size=10,
        color=SLATE,
        bold=True,
    )


def slide_2_problem(prs: Presentation) -> None:
    s = add_slide(prs)
    apply_chrome(s, 2)
    title_block(s, "The review-intelligence gap", eyebrow="Problem & motivation")

    # Left column: bullets
    add_bullets(
        s,
        Inches(0.6),
        Inches(2.4),
        Inches(6.8),
        Inches(4.2),
        [
            "Companies receive thousands of reviews daily across channels",
            "Manual analysis is slow, inconsistent, and expensive to scale",
            "Opportunity: specialised AI agents, each best-in-class at one job",
            "Goal: sentiment + competitive context, one pipeline",
        ],
        size=18,
        gap=1.5,
    )

    # Right column: big stat callout on a card
    add_round_rect(s, Inches(7.8), Inches(2.4), Inches(5.0), Inches(4.2), CARD)
    add_text(
        s,
        Inches(7.8),
        Inches(2.7),
        Inches(5.0),
        Inches(0.4),
        "REVIEWS PER LARGE BRAND / DAY",
        size=10,
        color=GOLD,
        bold=True,
        align="center",
    )
    add_text(
        s,
        Inches(7.8),
        Inches(3.25),
        Inches(5.0),
        Inches(1.7),
        "10,000+",
        font=H_FONT,
        size=72,
        color=WHITE,
        bold=True,
        align="center",
    )
    add_text(
        s,
        Inches(7.8),
        Inches(5.15),
        Inches(5.0),
        Inches(1.3),
        "Too many for a human team\nto read, tag, and act on in real time.",
        size=14,
        color=SLATE,
        italic=True,
        align="center",
        line_spacing=1.3,
    )


def slide_3_overview(prs: Presentation) -> None:
    s = add_slide(prs)
    apply_chrome(s, 3)
    title_block(s, "Pipeline at a glance", eyebrow="System overview")

    # Pipeline: 6 stages across the slide
    stages = [
        ("Review", "raw text"),
        ("Guardrail", "validate / clean"),
        ("Sentiment\nAgent", "BERT classifier"),
        ("HITL", "human check"),
        ("Market\nAgent", "web search"),
        ("Report\nAgent", "synthesise"),
    ]

    n = len(stages)
    left = Inches(0.6)
    top = Inches(2.5)
    total_w = Inches(12.1)
    box_w = Inches(1.75)
    box_h = Inches(1.55)
    gap = (total_w - Inches(1.75 * n)) / (n - 1) if n > 1 else Emu(0)

    for i, (label, caption) in enumerate(stages):
        x = left + i * (box_w + gap)
        # Highlight HITL in gold as a special checkpoint
        is_hitl = label == "HITL"
        fill = GOLD if is_hitl else CARD
        text_color = NAVY if is_hitl else WHITE
        cap_color = CARD_DARKER if is_hitl else SLATE

        add_round_rect(s, x, top, box_w, box_h, fill)
        add_text(
            s,
            x,
            top + Inches(0.25),
            box_w,
            Inches(0.8),
            label,
            font=H_FONT,
            size=14,
            color=text_color,
            bold=True,
            align="center",
        )
        add_text(
            s,
            x,
            top + Inches(1.0),
            box_w,
            Inches(0.45),
            caption,
            size=10,
            color=cap_color,
            align="center",
            italic=True,
        )

        # Gold arrow chevron between boxes
        if i < n - 1:
            chev_x = x + box_w + Inches(0.04)
            chev_y = top + box_h / 2 - Inches(0.12)
            chev = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_TRIANGLE,
                chev_x,
                chev_y,
                Inches(0.22),
                Inches(0.24),
            )
            solid_fill(chev, GOLD)

    # Caption below the pipeline
    add_text(
        s,
        Inches(0.6),
        Inches(4.6),
        Inches(12),
        Inches(0.4),
        "3 specialised agents  +  1 orchestrator     ·     2 tools:  BERT classifier  +  Web search",
        size=15,
        color=GOLD,
        bold=True,
        align="center",
    )

    # Bottom-row: labelled counts
    def stat(x, big, small):
        add_text(s, x, Inches(5.35), Inches(3.5), Inches(0.7), big,
                 font=H_FONT, size=44, color=WHITE, bold=True, align="center")
        add_text(s, x, Inches(6.1), Inches(3.5), Inches(0.4), small,
                 size=11, color=SLATE, bold=True, align="center")

    stat(Inches(1.4), "3", "SPECIALIST AGENTS")
    stat(Inches(4.9), "2", "TOOLS (BERT + SEARCH)")
    stat(Inches(8.4), "1", "ORCHESTRATOR (CREWAI)")


def slide_4_stack(prs: Presentation) -> None:
    s = add_slide(prs)
    apply_chrome(s, 4)
    title_block(s, "What the system is built on", eyebrow="Technical stack")

    rows = [
        ("Agent framework", "CrewAI 0.5.0",        "Simple sequential orchestration"),
        ("LLM backend",     "GPT-4o-mini",          "Cost-efficient, reliable JSON output"),
        ("DL framework",    "PyTorch + HuggingFace","Industry standard for BERT"),
        ("Search",          "DuckDuckGo",           "Free, no API key required"),
        ("Language",        "Python 3.10",          "Broad library compatibility"),
    ]

    left = Inches(0.6)
    top = Inches(2.35)
    w_comp = Inches(3.4)
    w_tech = Inches(3.9)
    w_why  = Inches(4.8)
    row_h = Inches(0.82)

    # Header row (gold underline, no cell fill)
    add_text(s, left, top, w_comp, Inches(0.35), "COMPONENT",
             size=11, color=GOLD, bold=True)
    add_text(s, left + w_comp, top, w_tech, Inches(0.35), "TECHNOLOGY",
             size=11, color=GOLD, bold=True)
    add_text(s, left + w_comp + w_tech, top, w_why, Inches(0.35), "WHY",
             size=11, color=GOLD, bold=True)
    add_rect(s, left, top + Inches(0.4), w_comp + w_tech + w_why, Inches(0.02), GOLD)

    # Body rows — alternating card fill for legibility
    row_top = top + Inches(0.55)
    for i, (comp, tech, why) in enumerate(rows):
        if i % 2 == 0:
            add_rect(s, left, row_top, w_comp + w_tech + w_why, row_h, CARD)
        pad_y = row_top + Inches(0.2)
        add_text(s, left + Inches(0.25), pad_y, w_comp - Inches(0.2), row_h,
                 comp, size=16, color=WHITE, bold=True)
        add_text(s, left + w_comp, pad_y, w_tech, row_h,
                 tech, font=H_FONT, size=16, color=GOLD, italic=True)
        add_text(s, left + w_comp + w_tech, pad_y, w_why - Inches(0.25), row_h,
                 why, size=14, color=SLATE)
        row_top += row_h


def slide_5_bert(prs: Presentation) -> None:
    s = add_slide(prs)
    apply_chrome(s, 5)
    title_block(s, "Fine-tuning BERT on Amazon Reviews", eyebrow="BERT sentiment model")

    # Left: bullets
    add_bullets(
        s,
        Inches(0.6),
        Inches(2.4),
        Inches(6.5),
        Inches(4.5),
        [
            "Dataset: Amazon Reviews All_Beauty — 10,000 rows",
            "3 classes: negative / neutral / positive",
            "Fine-tuned bert-base-uncased, 3 epochs",
            "Weighted F1 chosen because of class imbalance",
        ],
        size=17,
        gap=1.5,
    )

    # Right: metric cards
    def metric_card(x, y, big, label):
        add_round_rect(s, x, y, Inches(2.8), Inches(1.9), CARD)
        add_text(s, x, y + Inches(0.35), Inches(2.8), Inches(1.0),
                 big, font=H_FONT, size=44, color=GOLD, bold=True, align="center")
        add_text(s, x, y + Inches(1.35), Inches(2.8), Inches(0.4),
                 label, size=11, color=SLATE, bold=True, align="center")

    metric_card(Inches(7.4), Inches(2.4), "86.2%", "ACCURACY · EPOCH 2")
    metric_card(Inches(10.4), Inches(2.4), "0.859", "WEIGHTED F1 · BEST")

    # Class distribution bars (below cards)
    add_text(s, Inches(7.4), Inches(4.55), Inches(5.8), Inches(0.35),
             "CLASS DISTRIBUTION", size=10, color=GOLD, bold=True)

    def dist_bar(y, label, pct, color):
        bar_left = Inches(7.4)
        track_x = bar_left + Inches(1.3)
        track_w = Inches(2.8)            # track region only
        pct_x = bar_left + Inches(4.2)   # percentage label to the right of track
        filled = Inches(2.8 * pct / 100)
        # Tiny floor so 0-like classes render a sliver; preserves proportion.
        filled = max(filled, Inches(0.08))
        add_text(s, bar_left, y, Inches(1.3), Inches(0.3),
                 label, size=12, color=WHITE, bold=True)
        add_rect(s, track_x, y + Inches(0.05), track_w, Inches(0.2), CARD)
        add_rect(s, track_x, y + Inches(0.05), filled, Inches(0.2), color)
        add_text(s, pct_x, y, Inches(1.0), Inches(0.3),
                 f"{pct}%", size=12, color=SLATE, bold=True)

    dist_bar(Inches(5.0), "Negative",  14.4, RED)
    dist_bar(Inches(5.4), "Neutral",    9.8, SLATE)
    dist_bar(Inches(5.8), "Positive",  75.8, GREEN)

    # Footnote — tucked under the bullets on the LEFT, well clear of the footer
    add_text(
        s,
        Inches(0.6),
        Inches(6.3),
        Inches(6.5),
        Inches(0.5),
        "Early stopping triggered after Epoch 2 —\nlater epochs did not improve F1.",
        size=12,
        color=SLATE,
        italic=True,
        line_spacing=1.25,
    )


def slide_6_tool(prs: Presentation) -> None:
    s = add_slide(prs)
    apply_chrome(s, 6)
    title_block(s, "Wrapping the model as a CrewAI tool", eyebrow="BERT as a tool")

    # Left column bullets
    add_bullets(
        s,
        Inches(0.6),
        Inches(2.4),
        Inches(6.2),
        Inches(4.5),
        [
            "Tool name: sentiment_classifier",
            "Input:  review text (string)",
            "Output: JSON — label, confidence, per-class scores",
            "Lazy loading — model weights cached after first call",
            "Agent MUST call the tool — cannot guess sentiment",
        ],
        size=17,
        gap=1.5,
    )

    # Right column: mock JSON output on a code card
    card_x = Inches(7.1)
    card_y = Inches(2.4)
    card_w = Inches(5.7)
    card_h = Inches(4.1)

    add_round_rect(s, card_x, card_y, card_w, card_h, CARD_DARKER)
    # Faux terminal header bar
    add_rect(s, card_x, card_y, card_w, Inches(0.38), CARD)
    add_text(s, card_x + Inches(0.25), card_y + Inches(0.07),
             card_w - Inches(0.5), Inches(0.3),
             "sentiment_classifier → output",
             size=10, color=SLATE, bold=True)
    add_text(s, card_x + Inches(0.25), card_y + Inches(0.07),
             card_w - Inches(0.5), Inches(0.3),
             "JSON", size=10, color=GOLD, bold=True, align="right")

    code = (
        "{\n"
        '  "label": "positive",\n'
        '  "confidence": 0.94,\n'
        '  "scores": {\n'
        '    "negative": 0.02,\n'
        '    "neutral":  0.04,\n'
        '    "positive": 0.94\n'
        "  }\n"
        "}"
    )
    add_text(
        s,
        card_x + Inches(0.35),
        card_y + Inches(0.55),
        card_w - Inches(0.7),
        card_h - Inches(0.7),
        code,
        font="Consolas",
        size=15,
        color=WHITE,
        line_spacing=1.35,
    )

    # Caption under code card
    add_text(
        s,
        card_x,
        card_y + card_h + Inches(0.1),
        card_w,
        Inches(0.3),
        "Deterministic, structured — safe to pass to the next agent.",
        size=11,
        color=GOLD,
        italic=True,
        align="center",
    )


def slide_7_multiagent(prs: Presentation) -> None:
    s = add_slide(prs)
    apply_chrome(s, 7)
    title_block(s, "Three specialists, one orchestrator", eyebrow="Multi-agent orchestration")

    agents = [
        ("01", "Sentiment\nSpecialist",
         "Calls the BERT tool on the review and emits a structured sentiment verdict.",
         "TOOL: sentiment_classifier"),
        ("02", "Market Research\nSpecialist",
         "Queries DuckDuckGo for the product category and summarises competitive context.",
         "TOOL: web_search"),
        ("03", "Report\nSynthesiser",
         "Combines the sentiment verdict and market brief into a single intelligence report.",
         "NO TOOLS — reasoning only"),
    ]

    left = Inches(0.6)
    top = Inches(2.4)
    card_w = Inches(4.0)
    card_h = Inches(3.9)
    gap = Inches(0.17)

    for i, (num, name, desc, tool) in enumerate(agents):
        x = left + i * (card_w + gap)
        # Card
        add_round_rect(s, x, top, card_w, card_h, CARD)
        # Gold accent bar at top of card
        add_rect(s, x, top, card_w, Inches(0.08), GOLD)
        # Oversized number
        add_text(s, x + Inches(0.3), top + Inches(0.2), Inches(1.2), Inches(0.9),
                 num, font=H_FONT, size=48, color=GOLD, bold=True)
        # Agent name
        add_text(s, x + Inches(0.3), top + Inches(1.25), card_w - Inches(0.6),
                 Inches(1.1),
                 name, font=H_FONT, size=22, color=WHITE, bold=True,
                 line_spacing=1.15)
        # Description
        add_text(s, x + Inches(0.3), top + Inches(2.35), card_w - Inches(0.6),
                 Inches(1.1),
                 desc, size=13, color=SLATE, line_spacing=1.35)
        # Tool tag at bottom
        add_rect(s, x + Inches(0.3), top + card_h - Inches(0.6),
                 card_w - Inches(0.6), Inches(0.02), GOLD_DIM)
        add_text(s, x + Inches(0.3), top + card_h - Inches(0.5),
                 card_w - Inches(0.6), Inches(0.35),
                 tool, size=10, color=GOLD, bold=True)

    # Footnote
    add_text(
        s,
        Inches(0.6),
        Inches(6.55),
        Inches(12),
        Inches(0.4),
        "Sequential process via CrewAI Crew  ·  context passed between tasks",
        size=13,
        color=SLATE,
        italic=True,
        align="center",
    )


def slide_8_hitl(prs: Presentation) -> None:
    s = add_slide(prs)
    apply_chrome(s, 8)
    title_block(s, "Human-in-the-loop after sentiment", eyebrow="HITL checkpoint")

    # Left column: flow bullets
    add_bullets(
        s,
        Inches(0.6),
        Inches(2.4),
        Inches(6.5),
        Inches(4.5),
        [
            "Checkpoint fires after Task 1 (sentiment)",
            "Human sees the label + full confidence scores",
            "Approve (y) lets the pipeline continue",
            "Reject (n) raises ValueError — pipeline halts cleanly",
            "Why here: sentiment drives every downstream decision",
        ],
        size=16,
        gap=1.45,
    )

    # Right column: decision gate — two stacked cards
    gate_x = Inches(7.6)
    gate_w = Inches(5.2)

    # y card
    add_round_rect(s, gate_x, Inches(2.4), gate_w, Inches(2.05), CARD)
    add_rect(s, gate_x, Inches(2.4), Inches(0.1), Inches(2.05), GREEN)
    add_text(s, gate_x + Inches(0.5), Inches(2.55), Inches(1.2), Inches(0.9),
             "y", font=H_FONT, size=62, color=GREEN, bold=True)
    add_text(s, gate_x + Inches(1.7), Inches(2.7), gate_w - Inches(2.0),
             Inches(0.4),
             "APPROVE", size=11, color=GOLD, bold=True)
    add_text(s, gate_x + Inches(1.7), Inches(3.05), gate_w - Inches(2.0),
             Inches(1.2),
             "Proceed to Market Agent →\nReport Agent → final brief.",
             size=14, color=WHITE, line_spacing=1.3)

    # n card
    add_round_rect(s, gate_x, Inches(4.65), gate_w, Inches(2.05), CARD)
    add_rect(s, gate_x, Inches(4.65), Inches(0.1), Inches(2.05), RED)
    add_text(s, gate_x + Inches(0.5), Inches(4.8), Inches(1.2), Inches(0.9),
             "n", font=H_FONT, size=62, color=RED, bold=True)
    add_text(s, gate_x + Inches(1.7), Inches(4.95), gate_w - Inches(2.0),
             Inches(0.4),
             "REJECT", size=11, color=GOLD, bold=True)
    add_text(s, gate_x + Inches(1.7), Inches(5.3), gate_w - Inches(2.0),
             Inches(1.3),
             "ValueError raised.\nOrchestrator logs + exits cleanly —\nno hallucinated report.",
             size=14, color=WHITE, line_spacing=1.3)


def slide_9_eval(prs: Presentation) -> None:
    s = add_slide(prs)
    apply_chrome(s, 9)
    title_block(s, "Stress-testing the pipeline", eyebrow="Evaluation & robustness")

    # Left: big PASS stat card  (tightened — used to collide with pills row)
    card_x = Inches(0.6)
    card_y = Inches(2.4)
    card_h = Inches(3.8)
    add_round_rect(s, card_x, card_y, Inches(4.8), card_h, CARD)
    add_text(s, card_x, card_y + Inches(0.3), Inches(4.8), Inches(0.4),
             "EDGE-CASE REGRESSION", size=11, color=GOLD, bold=True, align="center")
    add_text(s, card_x, card_y + Inches(0.75), Inches(4.8), Inches(2.0),
             "6 / 6", font=H_FONT, size=110, color=GOLD, bold=True, align="center")
    add_text(s, card_x, card_y + Inches(2.75), Inches(4.8), Inches(0.5),
             "PASS", size=22, color=WHITE, bold=True, align="center")
    add_text(s, card_x, card_y + Inches(3.2), Inches(4.8), Inches(0.5),
             "Zero crashes across every hostile input.",
             size=12, color=SLATE, italic=True, align="center")

    # Right: enumerated cases as a 2x3 grid
    cases = [
        "Empty string",
        "600-word lorem-ipsum",
        "French (non-English)",
        "Abusive / offensive",
        "Numbers only",
        "Single word",
    ]

    grid_x = Inches(5.7)
    grid_y = Inches(2.4)
    cell_w = Inches(3.55)
    cell_h = Inches(0.85)
    col_gap = Inches(0.15)
    row_gap = Inches(0.15)

    for i, case in enumerate(cases):
        col = i % 2
        row = i // 2
        x = grid_x + col * (cell_w + col_gap)
        y = grid_y + row * (cell_h + row_gap)
        add_round_rect(s, x, y, cell_w, cell_h, CARD_DARKER)
        add_rect(s, x, y, Inches(0.08), cell_h, GOLD)
        add_text(s, x + Inches(0.3), y + Inches(0.2), Inches(0.5), Inches(0.55),
                 f"{i+1:02d}", font=H_FONT, size=20, color=GOLD, bold=True)
        add_text(s, x + Inches(1.0), y + Inches(0.27), cell_w - Inches(1.1),
                 Inches(0.5),
                 case, size=14, color=WHITE, bold=True)

    # Bottom strip: guardrail + logging pills — sits between cards and footer
    pills_top = Inches(6.4)
    pills = [
        ("GUARDRAIL", "min 3 / max 1000 words  ·  sanitised"),
        ("ERROR HANDLING", "try / except on every agent call"),
        ("AUDIT", "timestamped JSON run log"),
    ]
    pill_w = Inches(4.15)
    for i, (tag, body) in enumerate(pills):
        x = Inches(0.6) + i * (pill_w + Inches(0.1))
        add_text(s, x, pills_top, pill_w, Inches(0.28), tag,
                 size=9, color=GOLD, bold=True)
        add_text(s, x, pills_top + Inches(0.28), pill_w, Inches(0.35), body,
                 size=11, color=WHITE)


def slide_10_demo(prs: Presentation) -> None:
    s = add_slide(prs)
    apply_chrome(s, 10)
    # Eyebrow
    add_text(s, Inches(0.6), Inches(0.55), Inches(10), Inches(0.3),
             "LIVE DEMONSTRATION", size=11, color=GOLD, bold=True)

    # Big hero title, centered vertically
    add_text(
        s,
        Inches(0.6),
        Inches(1.9),
        Inches(12.1),
        Inches(2.0),
        "Demo time",
        font=H_FONT,
        size=120,
        color=WHITE,
        bold=True,
        align="center",
    )
    # Accent bar — clearly BELOW the title's descent (title box ends at 3.9)
    add_rect(s, Inches(6.0), Inches(4.1), Inches(1.3), Inches(0.07), GOLD)

    # 3 supporting lines below
    bullets_y = Inches(4.7)
    lines = [
        ("01", "Input", "3 product reviews — positive, negative, mixed"),
        ("02", "Checkpoint", "Watch the HITL gate fire in real time"),
        ("03", "Output", "Sentiment + confidence + market context"),
    ]
    for i, (num, tag, body) in enumerate(lines):
        y = bullets_y + i * Inches(0.7)
        add_text(s, Inches(3.2), y, Inches(0.7), Inches(0.55),
                 num, font=H_FONT, size=24, color=GOLD, bold=True)
        add_text(s, Inches(4.0), y + Inches(0.05), Inches(2.0), Inches(0.55),
                 tag, size=15, color=GOLD, bold=True)
        add_text(s, Inches(5.9), y + Inches(0.05), Inches(6.5), Inches(0.55),
                 body, size=16, color=WHITE)


def slide_11_conclusion(prs: Presentation) -> None:
    s = add_slide(prs)
    apply_chrome(s, 11)
    title_block(s, "What we achieved — and what's next", eyebrow="Conclusion")

    # Three columns: Achieved / Limitations / Future
    col_y = Inches(2.2)
    col_w = Inches(4.0)
    col_h = Inches(3.7)
    col_gap = Inches(0.17)
    col_x0 = Inches(0.6)

    columns = [
        ("ACHIEVED", GREEN, [
            "End-to-end multi-agent system",
            "BERT fine-tuned · 86.2% accuracy",
            "Human-in-the-loop checkpoint",
            "6 / 6 edge cases PASS",
        ]),
        ("LIMITATIONS", GOLD, [
            "Class imbalance — ~75% positive",
            "DuckDuckGo rate limits / reliability",
            "CrewAI 0.5.0 API constraints",
        ]),
        ("FUTURE WORK", SLATE, [
            "Migrate to LangGraph",
            "Streamlit UI for non-technical users",
            "Larger + more balanced dataset",
        ]),
    ]

    for i, (heading, accent, items) in enumerate(columns):
        x = col_x0 + i * (col_w + col_gap)
        add_round_rect(s, x, col_y, col_w, col_h, CARD)
        add_rect(s, x, col_y, col_w, Inches(0.08), accent)
        add_text(s, x + Inches(0.3), col_y + Inches(0.25), col_w - Inches(0.6),
                 Inches(0.4),
                 heading, size=13, color=accent, bold=True)
        add_bullets(
            s,
            x + Inches(0.3),
            col_y + Inches(0.85),
            col_w - Inches(0.6),
            col_h - Inches(1.1),
            items,
            size=14,
            gap=1.4,
            bullet_color=accent,
        )

    # Thank-you strip at the bottom — clear gap from cards above, footer below
    strip_y = Inches(6.25)
    add_rect(s, Inches(0.6), strip_y, Inches(12.17), Inches(0.05), GOLD)
    add_text(s, Inches(0.6), strip_y + Inches(0.2), Inches(7), Inches(0.55),
             "Thank you.",
             font=H_FONT, size=28, color=WHITE, bold=True)
    add_text(s, Inches(7.7), strip_y + Inches(0.3), Inches(5.07), Inches(0.55),
             "Questions?",
             font=H_FONT, size=22, color=GOLD, italic=True, align="right")


# ===========================================================================
# Assemble
# ===========================================================================
def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_1_title(prs)
    slide_2_problem(prs)
    slide_3_overview(prs)
    slide_4_stack(prs)
    slide_5_bert(prs)
    slide_6_tool(prs)
    slide_7_multiagent(prs)
    slide_8_hitl(prs)
    slide_9_eval(prs)
    slide_10_demo(prs)
    slide_11_conclusion(prs)

    prs.save(str(OUT))
    kb = OUT.stat().st_size / 1024
    print(f"Wrote {OUT}  ({kb:.1f} KB, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
